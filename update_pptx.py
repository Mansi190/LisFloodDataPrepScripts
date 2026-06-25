from pptx import Presentation
from pptx.util import Inches, Pt
import os

pptx_path = "ppts/LisFlood_Cold_Start_Presentation.pptx"
prs = Presentation(pptx_path)

# Find a blank layout or title & content layout
# Layout 5 is usually "Title Only" or 6 is "Blank"
# Let's try layout 5 (Title Only)
slide_layout = prs.slide_layouts[5]

# ---- Slide 6: Cold Start Plot ----
slide1 = prs.slides.add_slide(slide_layout)
shapes1 = slide1.shapes
title_shape1 = shapes1.title
if title_shape1:
    title_shape1.text = "Cold Start: River Discharge Hydrograph"

img_path1 = "discharge_plot.png"
if os.path.exists(img_path1):
    # Add image, centered
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    shapes1.add_picture(img_path1, left, top, width=width)

# ---- Slide 7: Warm Start Plot ----
slide2 = prs.slides.add_slide(slide_layout)
shapes2 = slide2.shapes
title_shape2 = shapes2.title
if title_shape2:
    title_shape2.text = "Warm Start: River Discharge Hydrograph"

img_path2 = "discharge_plot_warm.png"
if os.path.exists(img_path2):
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    shapes2.add_picture(img_path2, left, top, width=width)

# ---- Slide 8: Discussion Summary ----
# Layout 1 is usually "Title and Content"
slide_layout_content = prs.slide_layouts[1]
slide3 = prs.slides.add_slide(slide_layout_content)
title_shape3 = slide3.shapes.title
if title_shape3:
    title_shape3.text = "Hydrograph Comparison: Cold vs Warm Start"

body_shape = slide3.shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "Visual observations from the discharge plots:"

p = tf.add_paragraph()
p.text = "In the Cold Start scenario, the simulation begins with zero/bogus initial states. It takes a significant spin-up period for the stores to fill up, leading to an initially delayed or suppressed discharge curve."
p.level = 1

p2 = tf.add_paragraph()
p2.text = "In the Warm Start scenario, initial states are loaded from a previous end-of-run state. The model immediately produces physically realistic discharge outputs without the initial spin-up lag."
p2.level = 1

p3 = tf.add_paragraph()
p3.text = "Using a Warm Start drastically reduces computation time and yields accurate real-time forecasts instantly."
p3.level = 1

prs.save(pptx_path)
print("Presentation updated successfully!")
