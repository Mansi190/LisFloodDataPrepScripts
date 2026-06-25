"""
Build a polished LISFLOOD Cold-Start vs Warm-Start presentation.
LIGHT THEME variant — clean white backgrounds with coloured accents.

Data sourced from: cold.xml, warm_start.xml, output_dataset/
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── LIGHT colour palette ──────────────────────────────────────────
BG_WHITE      = RGBColor(0xF8, 0xFA, 0xFC)   # slide background
BG_CARD       = RGBColor(0xFF, 0xFF, 0xFF)   # card / panel bg (pure white)
BG_CARD_ALT   = RGBColor(0xF1, 0xF5, 0xF9)   # alternate light card
ACCENT_BLUE   = RGBColor(0x1D, 0x4E, 0xD8)   # deep blue
ACCENT_GREEN  = RGBColor(0x05, 0x96, 0x69)   # teal green
ACCENT_AMBER  = RGBColor(0xD9, 0x77, 0x06)   # warm amber
ACCENT_PURPLE = RGBColor(0x7C, 0x3A, 0xED)   # vibrant purple
ACCENT_ROSE   = RGBColor(0xE1, 0x1D, 0x48)   # rich rose
TEXT_DARK     = RGBColor(0x0F, 0x17, 0x2A)   # headings / primary text
TEXT_BODY     = RGBColor(0x33, 0x41, 0x55)   # body text
TEXT_MUTED    = RGBColor(0x64, 0x74, 0x8B)   # secondary / muted text
BORDER_LIGHT  = RGBColor(0xE2, 0xE8, 0xF0)   # subtle borders

# Table colours (light)
TABLE_HDR     = RGBColor(0x1E, 0x40, 0x6E)   # dark header row
TABLE_HDR_TXT = RGBColor(0xFF, 0xFF, 0xFF)   # white text on header
TABLE_ROW1    = RGBColor(0xFF, 0xFF, 0xFF)   # white rows
TABLE_ROW2    = RGBColor(0xF1, 0xF5, 0xF9)   # light gray alternating
TABLE_BODY_TXT = RGBColor(0x1E, 0x29, 0x3B)  # dark text in table body

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── helpers ────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if border_color:
        shp.line.color.rgb = border_color
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    try:
        shp.adjustments[0] = 0.04
    except Exception:
        pass
    # Add subtle shadow
    try:
        shp.shadow.inherit = False
    except Exception:
        pass
    return shp


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=14, color=TEXT_DARK, bold=False,
                  space_before=Pt(4), space_after=Pt(2), level=0,
                  font_name="Segoe UI"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    return p


def add_accent_line(slide, left, top, width, color):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def add_stat_card(slide, left, top, width, height, value, label, accent):
    card = add_rect(slide, left, top, width, height, BG_CARD, border_color=BORDER_LIGHT)
    add_accent_line(slide, left + Inches(0.15), top + Inches(0.08),
                    width - Inches(0.3), accent)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.25),
                 width - Inches(0.4), Inches(0.5),
                 value, font_size=22, color=accent, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.75),
                 width - Inches(0.4), Inches(0.4),
                 label, font_size=11, color=TEXT_MUTED,
                 alignment=PP_ALIGN.CENTER)


def styled_table(slide, left, top, width, rows_data, col_widths,
                 header_color=TABLE_HDR):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, left, top, width, Inches(0.35 * n_rows))
    tbl = table_shape.table

    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.name = "Segoe UI"
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if ri == 0:
                p.font.bold = True
                p.font.color.rgb = TABLE_HDR_TXT
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                p.font.color.rgb = TABLE_BODY_TXT
                cell.fill.solid()
                cell.fill.fore_color.rgb = (TABLE_ROW1 if ri % 2 == 1
                                             else TABLE_ROW2)
    return table_shape


# ── build presentation ─────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

script_dir = os.path.dirname(os.path.abspath(__file__))

# ======================================================================
# SLIDE 1 — Title
# ======================================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_WHITE)

# Decorative soft circle
circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                          Inches(10), Inches(-1.5), Inches(5), Inches(5))
circ.fill.solid()
circ.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
circ.line.fill.background()

add_accent_line(s, Inches(1.2), Inches(2.2), Inches(3), ACCENT_BLUE)
add_text_box(s, Inches(1.2), Inches(2.45), Inches(10), Inches(1),
             "LISFLOOD HYDROLOGICAL MODEL",
             font_size=38, color=TEXT_DARK, bold=True)
add_text_box(s, Inches(1.2), Inches(3.5), Inches(10), Inches(0.7),
             "Cold Start & Warm Start — Execution, Outputs & Analysis",
             font_size=20, color=ACCENT_BLUE)
add_text_box(s, Inches(1.2), Inches(4.4), Inches(9), Inches(0.6),
             "A comprehensive walkthrough of simulation setup, Docker execution, "
             "state-file transitions, output datasets, and discharge hydrograph comparison.",
             font_size=13, color=TEXT_MUTED)

add_rect(s, Inches(0), Inches(6.7), SLIDE_W, Inches(0.8), BG_CARD_ALT)
add_text_box(s, Inches(1.2), Inches(6.85), Inches(10), Inches(0.4),
             "Watershed Area: 781.836 km²   •   Daily Timestep (86 400 s)   "
             "•   Channel Sub-Step: 3 600 s   •   UTM Zone 43 / WGS 84   "
             "•   Docker: jrce1/lisflood",
             font_size=11, color=TEXT_MUTED)


# ======================================================================
# SLIDE 2 — What is Cold Start / Warm Start?
# ======================================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_WHITE)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "UNDERSTANDING THE TWO SIMULATION MODES",
             font_size=12, color=ACCENT_BLUE, bold=True)
add_accent_line(s, Inches(0.8), Inches(0.9), Inches(4), ACCENT_BLUE)

# — Cold card
add_rect(s, Inches(0.8), Inches(1.2), Inches(5.6), Inches(5.8),
         BG_CARD, border_color=BORDER_LIGHT)
add_accent_line(s, Inches(0.8), Inches(1.2), Inches(5.6), ACCENT_BLUE)
add_text_box(s, Inches(1.2), Inches(1.5), Inches(4.8), Inches(0.45),
             "❄  COLD START  (Spin-Up Run)",
             font_size=22, color=ACCENT_BLUE, bold=True)
tf = add_text_box(s, Inches(1.2), Inches(2.15), Inches(4.8), Inches(4.7),
                  "", font_size=12, color=TEXT_DARK)
for t, b in [
    ("What it is:",
     "The simulation starts from scratch — all internal water "
     "stores are empty or set to bogus placeholder values (0 or −9999)."),
    ("Purpose:",
     "Run a multi-year spin-up period so that soil moisture, "
     "groundwater zones, channel flows, snow cover and frost "
     "indices build up to physically realistic levels."),
    ("Key Point:",
     "Outputs from the early portion of the run are unreliable "
     "and typically discarded; only the end-of-run state maps "
     "(.end.nc) matter — they become the 'memory' of the watershed."),
    ("InitLisflood:",
     "Set to 1 (ON) — tells the model this is an initialization run."),
    ("Config File:", "cold.xml"),
]:
    add_paragraph(tf, f"▸ {t}", font_size=12, color=ACCENT_BLUE,
                  bold=True, space_before=Pt(10))
    add_paragraph(tf, f"   {b}", font_size=11, color=TEXT_BODY,
                  space_before=Pt(2))

# — Warm card
add_rect(s, Inches(6.9), Inches(1.2), Inches(5.6), Inches(5.8),
         BG_CARD, border_color=BORDER_LIGHT)
add_accent_line(s, Inches(6.9), Inches(1.2), Inches(5.6), ACCENT_GREEN)
add_text_box(s, Inches(7.3), Inches(1.5), Inches(4.8), Inches(0.45),
             "🔥  WARM START  (Operational Run)",
             font_size=22, color=ACCENT_GREEN, bold=True)
tf = add_text_box(s, Inches(7.3), Inches(2.15), Inches(4.8), Inches(4.7),
                  "", font_size=12, color=TEXT_DARK)
for t, b in [
    ("What it is:",
     "The simulation loads pre-stabilized state files (*.end.nc) "
     "produced at the end of the Cold Start, bypassing the spin-up."),
    ("Purpose:",
     "Immediately generate physically accurate discharge and "
     "water-balance outputs — ideal for real-time flood forecasting "
     "with daily weather feeds."),
    ("Key Point:",
     "Model outputs are trustworthy from timestep 1 because the "
     "'memory' of the watershed is already baked into the initial states."),
    ("InitLisflood:",
     "Set to 0 (OFF) — this is a real simulation, not initialization."),
    ("Config File:", "warm_start.xml"),
]:
    add_paragraph(tf, f"▸ {t}", font_size=12, color=ACCENT_GREEN,
                  bold=True, space_before=Pt(10))
    add_paragraph(tf, f"   {b}", font_size=11, color=TEXT_BODY,
                  space_before=Pt(2))


# ======================================================================
# SLIDE 3 — Cold Start: Configuration & Stats
# ======================================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_WHITE)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "COLD START — SIMULATION CONFIGURATION",
             font_size=12, color=ACCENT_BLUE, bold=True)
add_accent_line(s, Inches(0.8), Inches(0.9), Inches(4), ACCENT_BLUE)

# stat cards
add_stat_card(s, Inches(0.8),  Inches(1.2), Inches(2.8), Inches(1.15),
              "01/01/2024", "Simulation Start (StepStart)", ACCENT_BLUE)
add_stat_card(s, Inches(3.8),  Inches(1.2), Inches(2.8), Inches(1.15),
              "30/11/2024", "Simulation End (StepEnd)", ACCENT_BLUE)
add_stat_card(s, Inches(6.8),  Inches(1.2), Inches(2.8), Inches(1.15),
              "~336 Days", "Total Timesteps", ACCENT_AMBER)
add_stat_card(s, Inches(9.8),  Inches(1.2), Inches(2.8), Inches(1.15),
              "~50 seconds", "Run Time", ACCENT_GREEN)

# Docker command card
add_rect(s, Inches(0.8), Inches(2.65), Inches(11.7), Inches(1.1),
         BG_CARD_ALT, border_color=BORDER_LIGHT)
add_text_box(s, Inches(1.1), Inches(2.75), Inches(3), Inches(0.35),
             "DOCKER COMMAND", font_size=10, color=ACCENT_AMBER, bold=True)
add_text_box(s, Inches(1.1), Inches(3.1), Inches(11), Inches(0.5),
             'docker run --rm -v '
             '"filepath/LisFloodDataPrepScripts":/input  '
             'jrce1/lisflood  /input/cold.xml',
             font_size=13, color=ACCENT_BLUE, font_name="Consolas")

# Init state table
add_text_box(s, Inches(0.8), Inches(3.95), Inches(6), Inches(0.3),
             "INITIAL STATE VALUES  (no prior knowledge)",
             font_size=11, color=ACCENT_BLUE, bold=True)

state_rows = [
    ("Variable", "Init Value", "Description"),
    ("LZInitValue", "−9999", "Lower groundwater zone (steady-state)"),
    ("UZInitValue", "0", "Upper groundwater zone storage"),
    ("ThetaInit 1/2/3", "−9999", "Soil moisture layers (field capacity)"),
    ("ChanCrossSectionInit", "−9999", "Channel cross-section (½ bankfull)"),
    ("SnowCover A/B/C", "0", "Snow cover in elevation zones"),
    ("FrostIndexInitValue", "0", "Ground frost / freeze-thaw index"),
    ("CumIntInitValue", "0", "Cumulative rainfall interception"),
    ("DSLRInitValue", "1", "Days since last rain"),
]
styled_table(s, Inches(0.8), Inches(4.25), Inches(7.5), state_rows,
             [Inches(2.5), Inches(1.2), Inches(3.8)])

# Outputs enabled card
add_rect(s, Inches(8.8), Inches(3.95), Inches(3.7), Inches(3.2),
         BG_CARD, border_color=BORDER_LIGHT)
add_text_box(s, Inches(9.1), Inches(4.05), Inches(3.2), Inches(0.3),
             "OUTPUTS ENABLED",
             font_size=11, color=ACCENT_GREEN, bold=True)

tf = add_text_box(s, Inches(9.1), Inches(4.35), Inches(3.2), Inches(2.7),
                  "", font_size=12, color=TEXT_DARK)
for chk, desc in [
    ("✅ repEndMaps = 1",      "End-of-run state maps"),
    ("✅ repStateMaps = 1",    "State maps at report steps"),
    ("✅ repDischargeTs = 1",  "Discharge time series (.tss)"),
    ("✅ repDischargeMaps = 1","Discharge spatial maps (.nc)"),
    ("✅ InitLisflood = 1",    "Cold initialization mode ON"),
]:
    add_paragraph(tf, chk, font_size=10,
                  color=ACCENT_GREEN, bold=True, space_before=Pt(5))
    add_paragraph(tf, f"  {desc}", font_size=10, color=TEXT_MUTED)


# ======================================================================
# SLIDE 4 — Cold Start Discharge Plot
# ======================================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_WHITE)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "COLD START — DISCHARGE HYDROGRAPH",
             font_size=12, color=ACCENT_BLUE, bold=True)
add_accent_line(s, Inches(0.8), Inches(0.9), Inches(4), ACCENT_BLUE)

add_text_box(s, Inches(0.8), Inches(1.05), Inches(11), Inches(0.45),
             "Simulated river discharge (m³/s) at the reporting station "
             "over ~336 daily timesteps   •   01/01/2024 → 30/11/2024",
             font_size=12, color=TEXT_MUTED)

add_rect(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.9),
         BG_CARD, border_color=BORDER_LIGHT)
img = os.path.join(script_dir, "discharge_plot.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(1.0), Inches(1.75),
                         width=Inches(11.3), height=Inches(4.6))

add_text_box(s, Inches(1.0), Inches(6.65), Inches(11), Inches(0.5),
             "▸ Notice the initial flat period (~0 m³/s) — the spin-up "
             "phase where internal stores are still filling up from empty.",
             font_size=11, color=ACCENT_AMBER)


# ======================================================================
# SLIDE 5 — Warm Start: Configuration & Stats
# ======================================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_WHITE)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "WARM START — SIMULATION CONFIGURATION",
             font_size=12, color=ACCENT_GREEN, bold=True)
add_accent_line(s, Inches(0.8), Inches(0.9), Inches(4), ACCENT_GREEN)

# stat cards
add_stat_card(s, Inches(0.8),  Inches(1.2), Inches(2.8), Inches(1.15),
              "01/12/2024", "Simulation Start (StepStart)", ACCENT_GREEN)
add_stat_card(s, Inches(3.8),  Inches(1.2), Inches(2.8), Inches(1.15),
              "31/12/2024", "Simulation End (StepEnd)", ACCENT_GREEN)
add_stat_card(s, Inches(6.8),  Inches(1.2), Inches(2.8), Inches(1.15),
              "~31 Days", "Total Timesteps (Dec 2024)", ACCENT_AMBER)
add_stat_card(s, Inches(9.8),  Inches(1.2), Inches(2.8), Inches(1.15),
              "~10 seconds", "Run Time (est.)", ACCENT_GREEN)

# Docker command card
add_rect(s, Inches(0.8), Inches(2.65), Inches(11.7), Inches(1.1),
         BG_CARD_ALT, border_color=BORDER_LIGHT)
add_text_box(s, Inches(1.1), Inches(2.75), Inches(3), Inches(0.35),
             "DOCKER COMMAND", font_size=10, color=ACCENT_AMBER, bold=True)
add_text_box(s, Inches(1.1), Inches(3.1), Inches(11), Inches(0.5),
             'docker run --rm -v '
             '"filepath/LisFloodDataPrepScripts":/input  '
             'jrce1/lisflood  /input/warm_start.xml',
             font_size=13, color=ACCENT_GREEN, font_name="Consolas")

# State files loaded from cold start
add_text_box(s, Inches(0.8), Inches(4.0), Inches(8), Inches(0.35),
             "INITIAL STATES LOADED FROM COLD START  (*.end.nc files)",
             font_size=11, color=ACCENT_GREEN, bold=True)

wf_rows = [
    ("File", "Variable", "Description"),
    ("frost.end.nc", "FrostIndexInitValue",
     "Ground frost / freeze-thaw index"),
    ("cumi.end.nc", "CumIntInitValue",
     "Cumulative rainfall interception"),
    ("dslr.end.nc", "DSLRInitValue",
     "Days since last rain event"),
    ("cumf.end.nc", "CumIntForestInitValue",
     "Cumulative interception (forest)"),
    ("dslf.end.nc", "DSLRForestInitValue",
     "Days since last rain (forest)"),
    ("dsli.end.nc", "DSLRIrrigationInitValue",
     "Days since last rain (irrigated)"),
    ("cseal.end.nc", "CumIntSealedInitValue",
     "Cumulative interception (sealed)"),
]
styled_table(s, Inches(0.8), Inches(4.25), Inches(11.7), wf_rows,
             [Inches(2.2), Inches(3.5), Inches(6.0)])

add_text_box(s, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
             "PathInit = /input/output_dataset/out  (reads cold start's "
             "output)     PathOut = /input/output_dataset/out_warm  "
             "(writes to a separate directory)",
             font_size=10, color=TEXT_MUTED)


# ======================================================================
# SLIDE 6 — Warm Start Discharge Plot
# ======================================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_WHITE)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "WARM START — DISCHARGE HYDROGRAPH",
             font_size=12, color=ACCENT_GREEN, bold=True)
add_accent_line(s, Inches(0.8), Inches(0.9), Inches(4), ACCENT_GREEN)

add_text_box(s, Inches(0.8), Inches(1.05), Inches(11), Inches(0.45),
             "Simulated river discharge (m³/s) at the reporting station "
             "over ~31 daily timesteps   •   01/12/2024 → 31/12/2024",
             font_size=12, color=TEXT_MUTED)

add_rect(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.9),
         BG_CARD, border_color=BORDER_LIGHT)
img = os.path.join(script_dir, "discharge_plot_warm.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(1.0), Inches(1.75),
                         width=Inches(11.3), height=Inches(4.6))

add_text_box(s, Inches(1.0), Inches(6.65), Inches(11), Inches(0.5),
             "▸ The warm start produces realistic discharge immediately "
             "— notice the rainfall-driven peak around day 342-343 "
             "(~0.038 m³/s), consistent with a real storm event.",
             font_size=11, color=ACCENT_AMBER)


# ======================================================================
# SLIDE 7 — Output Files Summary
# ======================================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_WHITE)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "OUTPUT FILES — SUMMARY & DESCRIPTION",
             font_size=12, color=ACCENT_PURPLE, bold=True)
add_accent_line(s, Inches(0.8), Inches(0.9), Inches(4), ACCENT_PURPLE)

# — Left card: End-state maps
add_rect(s, Inches(0.8), Inches(1.2), Inches(5.9), Inches(5.8),
         BG_CARD, border_color=BORDER_LIGHT)
add_accent_line(s, Inches(0.8), Inches(1.2), Inches(5.9), ACCENT_BLUE)

add_text_box(s, Inches(1.1), Inches(1.4), Inches(5.4), Inches(0.4),
             "END-OF-RUN STATE MAPS  (.end.nc)",
             font_size=14, color=ACCENT_BLUE, bold=True)
add_text_box(s, Inches(1.1), Inches(1.85), Inches(5.4), Inches(0.5),
             "NetCDF snapshots of every grid cell at the final timestep. "
             "Cold start produces 30 .end.nc files — the watershed's "
             "'memory' for warm-start initialization.",
             font_size=11, color=TEXT_BODY)

end_rows = [
    ("Category", "Files", "Description"),
    ("Channel", "chanq, chcro .end.nc", "Discharge & cross-section"),
    ("Snow", "scova/b/c .end.nc", "Snow cover per elevation zone"),
    ("Soil Moisture", "tha/b/c .end.nc", "Soil moisture layers a/b/c"),
    ("Groundwater", "lz, uz, uzf, uzi .end.nc", "Lower/upper zone storage"),
    ("Interception", "cumi, cumf, cseal .end.nc", "Rainfall interception"),
    ("Days Since Rain", "dslr, dslf, dsli .end.nc", "General / forest / irrigated"),
    ("Frost", "frost .end.nc", "Ground frost index"),
    ("Overland Flow", "ofdir, offor, ofoth .end.nc", "Direct / forest / other"),
]
styled_table(s, Inches(1.0), Inches(2.45), Inches(5.5), end_rows,
             [Inches(1.3), Inches(1.9), Inches(2.3)])

# — Right card: Time Series + Spatial Maps
add_rect(s, Inches(7.1), Inches(1.2), Inches(5.5), Inches(5.8),
         BG_CARD, border_color=BORDER_LIGHT)
add_accent_line(s, Inches(7.1), Inches(1.2), Inches(5.5), ACCENT_GREEN)

add_text_box(s, Inches(7.4), Inches(1.4), Inches(5), Inches(0.4),
             "TIME SERIES & SPATIAL OUTPUTS",
             font_size=14, color=ACCENT_GREEN, bold=True)

# TSS table
add_text_box(s, Inches(7.4), Inches(1.85), Inches(5), Inches(0.3),
             "Time Series (.tss files)",
             font_size=11, color=ACCENT_AMBER, bold=True)
tss_rows = [
    ("File", "Source", "Description"),
    ("dis.tss", "Both", "Discharge at gauge (m³/s)"),
    ("chanqWin.tss", "Both", "Channel Q last sub-step"),
    ("rainUps.tss", "Warm", "Upstream rainfall"),
    ("tAvgUps.tss", "Warm", "Upstream avg temperature"),
    ("etUps.tss", "Warm", "Upstream evapotranspiration"),
    ("snowUps.tss", "Warm", "Upstream snowmelt"),
]
styled_table(s, Inches(7.3), Inches(2.2), Inches(5.1), tss_rows,
             [Inches(1.5), Inches(0.7), Inches(2.9)])

# Spatial maps
add_text_box(s, Inches(7.4), Inches(4.6), Inches(5), Inches(0.3),
             "Spatial Maps (.nc)",
             font_size=11, color=ACCENT_AMBER, bold=True)
nc_rows = [
    ("File", "Size", "Description"),
    ("dis.nc", "5.65 MB", "Spatially-distributed discharge"),
    ("chanq.nc, chcro.nc", "20-31 KB", "Channel state maps"),
    ("tha/b/c.nc", "~76 KB", "Soil moisture maps"),
    ("lz.nc, uz.nc", "76-80 KB", "Groundwater zone maps"),
]
styled_table(s, Inches(7.3), Inches(4.95), Inches(5.1), nc_rows,
             [Inches(1.8), Inches(0.8), Inches(2.5)])


# ======================================================================
# SLIDE 8 — Side-by-side comparison table
# ======================================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_WHITE)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "COLD START vs WARM START — AT A GLANCE",
             font_size=12, color=ACCENT_AMBER, bold=True)
add_accent_line(s, Inches(0.8), Inches(0.9), Inches(4), ACCENT_AMBER)

cmp = [
    ("Parameter", "Cold Start", "Warm Start"),
    ("Config File", "cold.xml", "warm_start.xml"),
    ("Simulation Period",
     "01/01/2024 → 30/11/2024",
     "01/12/2024 → 31/12/2024"),
    ("Duration", "~336 days (~11 months)",
     "~31 days (December 2024)"),
    ("Approx. Run Time", "~50 seconds", "~10 seconds"),
    ("Initial States",
     "All set to 0 or −9999 (empty watershed)",
     "7 × .end.nc files loaded from cold run"),
    ("InitLisflood", "1  (initialization mode ON)",
     "0  (operational mode)"),
    ("Spin-Up Required",
     "Yes — early output unreliable",
     "No — accurate from timestep 1"),
    ("Output Directory", "/output_dataset/out",
     "/output_dataset/out_warm"),
    ("Key Outputs",
     "30 .end.nc state maps + dis.tss",
     "33 .end.nc + 8 .tss + 25 spatial .nc"),
    ("Docker Image", "jrce1/lisflood", "jrce1/lisflood"),
    ("Use Case",
     "Build initial states / calibrate",
     "Real-time operational forecasting"),
]
styled_table(s, Inches(0.8), Inches(1.15), Inches(11.7), cmp,
             [Inches(2.5), Inches(4.6), Inches(4.6)])


# ======================================================================
# SLIDE 9 — Cold → Warm Workflow Diagram
# ======================================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_WHITE)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "COLD → WARM START WORKFLOW",
             font_size=12, color=ACCENT_PURPLE, bold=True)
add_accent_line(s, Inches(0.8), Inches(0.9), Inches(4), ACCENT_PURPLE)

# Step 1 card
add_rect(s, Inches(0.8), Inches(1.4), Inches(3.5), Inches(5.2),
         BG_CARD, border_color=BORDER_LIGHT)
add_accent_line(s, Inches(0.8), Inches(1.4), Inches(3.5), ACCENT_BLUE)
add_text_box(s, Inches(1.1), Inches(1.6), Inches(3), Inches(0.4),
             "STEP 1 — COLD START",
             font_size=16, color=ACCENT_BLUE, bold=True)
tf = add_text_box(s, Inches(1.1), Inches(2.15), Inches(3), Inches(4.2),
                  "", font_size=11, color=TEXT_BODY)
for line in [
    "▸ Run cold.xml for ~336 days",
    "▸ All states start at 0 / −9999",
    "▸ Watershed fills up over time",
    "▸ End-of-run: 30 state maps saved",
    "  as *.end.nc in /out/",
    "",
    "▸ Output: dis.tss + dis.nc",
    "  (discharge data for analysis)",
    "",
    "▸ Total storage: ~130 MB",
]:
    add_paragraph(tf, line, font_size=11, color=TEXT_BODY,
                  space_before=Pt(3))

# Arrow
add_text_box(s, Inches(4.5), Inches(3.3), Inches(1.6), Inches(1.5),
             "→", font_size=60, color=ACCENT_AMBER, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(4.3), Inches(4.3), Inches(2), Inches(0.8),
             "7 × .end.nc\nfiles transfer",
             font_size=11, color=ACCENT_AMBER, alignment=PP_ALIGN.CENTER)

# Step 2 card
add_rect(s, Inches(6.3), Inches(1.4), Inches(3.5), Inches(5.2),
         BG_CARD, border_color=BORDER_LIGHT)
add_accent_line(s, Inches(6.3), Inches(1.4), Inches(3.5), ACCENT_GREEN)
add_text_box(s, Inches(6.6), Inches(1.6), Inches(3), Inches(0.4),
             "STEP 2 — WARM START",
             font_size=16, color=ACCENT_GREEN, bold=True)
tf = add_text_box(s, Inches(6.6), Inches(2.15), Inches(3), Inches(4.2),
                  "", font_size=11, color=TEXT_BODY)
for line in [
    "▸ Run warm_start.xml for 31 days",
    "▸ 7 key states loaded from .end.nc",
    "▸ Remaining states use defaults",
    "▸ Physically realistic from day 1",
    "",
    "▸ Outputs written to /out_warm/",
    "▸ 8 .tss files + 25 spatial .nc",
    "  + 33 new .end.nc maps",
    "",
    "▸ Ready for next warm cycle!",
]:
    add_paragraph(tf, line, font_size=11, color=TEXT_BODY,
                  space_before=Pt(3))

# Result card
add_rect(s, Inches(10.2), Inches(1.4), Inches(2.5), Inches(5.2),
         BG_CARD, border_color=BORDER_LIGHT)
add_accent_line(s, Inches(10.2), Inches(1.4), Inches(2.5), ACCENT_ROSE)
add_text_box(s, Inches(10.4), Inches(1.6), Inches(2.1), Inches(0.4),
             "RESULTS",
             font_size=16, color=ACCENT_ROSE, bold=True)
tf = add_text_box(s, Inches(10.4), Inches(2.1), Inches(2.1), Inches(4.4),
                  "", font_size=11, color=TEXT_BODY)
result_items = [
    ("Cold Run", "30 .end.nc + 5 .tss\n+ dis.nc (5.65 MB)"),
    ("Warm Run", "33 .end.nc + 8 .tss\n+ 25 spatial .nc maps"),
    ("Peak Q (Cold)", "~1.19 m³/s\naround day 275"),
    ("Peak Q (Warm)", "~0.038 m³/s\naround day 342"),
    ("Run Time ↓", "50 s → 10 s\n(~80% faster)"),
]
for label, val in result_items:
    add_paragraph(tf, label, font_size=10, color=ACCENT_ROSE,
                  bold=True, space_before=Pt(8))
    add_paragraph(tf, val, font_size=10, color=TEXT_BODY,
                  space_before=Pt(1))


# ======================================================================
# SLIDE 10 — Thank You
# ======================================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_WHITE)

circ2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                           Inches(-2), Inches(4), Inches(6), Inches(6))
circ2.fill.solid()
circ2.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
circ2.line.fill.background()

add_accent_line(s, Inches(4.5), Inches(2.6), Inches(4.3), ACCENT_BLUE)
add_text_box(s, Inches(1), Inches(2.8), Inches(11.3), Inches(1),
             "Thank You", font_size=44, color=TEXT_DARK, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(1), Inches(3.9), Inches(11.3), Inches(0.6),
             "LISFLOOD  •  Cold Start & Warm Start Analysis",
             font_size=16, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(1), Inches(4.7), Inches(11.3), Inches(0.5),
             "For questions or further discussion, please reach out.",
             font_size=13, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ── save ───────────────────────────────────────────────────────────
out = os.path.join(script_dir, "ppts",
                   "LisFlood_Cold_Warm_Start_Presentation.pptx")
prs.save(out)
print(f"✅  Presentation saved → {out}")
print(f"   Total slides: {len(prs.slides)}")
